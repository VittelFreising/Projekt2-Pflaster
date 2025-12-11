import os
from pptx import Presentation
from pptx.util import Inches, Pt

def create_german_lab_ppt_no_spectra(image_folder, output_filename="Laborbericht_Pflaster_Analyse.pptx"):
    prs = Presentation()

    # --- 1. Titelfolie (封面) ---
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Materialcharakterisierung von Wundschnellverbänden"
    subtitle.text = "Mikroskopie, FTIR-Analyse & Mechanische Prüfung"

    # --- 数据列表 ---
    # 格式: [幻灯片标题, 文本内容, 图片文件名]
    slides_data = [
        # --- Slide 2: Trägermaterial (基材) ---
        [
            "1. Trägermaterial (Backing)", 
            "Mikroskopische Beobachtung:\n"
            "• Regelmäßige Wabenstruktur (Honeycomb-Struktur) zur Gewährleistung der Atmungsaktivität.\n"
            "• Sichtbare Perforation der Folie.\n\n"
            "Chemische Analyse (FTIR):\n"
            "• Identifiziert als: Polyethylen (PE).\n"
            "• Übereinstimmung (Hit Quality): >94%.\n"
            "• Nachweis: Charakteristische C-H Streckschwingungen (2915, 2848 cm⁻¹) und Biegeschwingungen (1460 cm⁻¹).",
            "IMG-20251208-WA0007.jpg"  # 您的蜂窝状图
        ],
        
        # --- Slide 3: Wundauflage (吸水垫) ---
        [
            "2. Wundauflage (Wound Pad)", 
            "Mikroskopische Beobachtung:\n"
            "• Ungeordnete Faserstruktur (Vliesstoff / Non-woven).\n"
            "• Hohe Porosität begünstigt die Aufnahme von Wundexsudat.\n\n"
            "Chemische Analyse (FTIR):\n"
            "• Identifiziert als: Zellulose (Cellulose).\n"
            "• Übereinstimmung (Hit Quality): ~90%.\n"
            "• Nachweis: Breite O-H Bande (3300-3400 cm⁻¹) und C-O Fingerprint-Region typisch für Polysaccharide.",
            "IMG-20251208-WA0006.jpg"  # 您的纤维图
        ],
        
        # --- Slide 4: Klebstoff (胶粘剂) ---
        [
            "3. Klebstoff & Liner (Adhesive)", 
            "Mikroskopische Beobachtung:\n"
            "• Inhomogener Auftrag des Klebstoffs.\n"
            "• Sichtbare Defekte: Blasen (Bubbles) und ringförmige Trocknungsstrukturen.\n\n"
            "Chemische Analyse (FTIR):\n"
            "• Klebstoff: Polyacrylat / Acrylat-Copolymer (z.B. Polybutylacrylat, Match ~92%).\n"
            "• Schutzfolie (Liner): Polydimethylsiloxan (Silikon/PDMS) als Trennschicht (Match ~90%).",
            "IMG-20251208-WA0009.jpg"  # 您的胶水气泡图
        ],

        # --- Slide 5: Mechanik (拉伸) ---
        [
            "4. Mechanische Eigenschaften (Zugversuch)", 
            "Versuchsaufbau:\n"
            "• Einachsiger Zugversuch bis zum Bruch.\n\n"
            "Beobachtung:\n"
            "• Deutliche Einschnürung (Necking) der Probe während der Belastung.\n"
            "• Weißbruch (Stress Whitening) im verformten Bereich.\n\n"
            "Fazit:\n"
            "• Das Material zeigt ein stark duktiles (plastisches) Verformungsverhalten, charakteristisch für Polyethylen.",
            "IMG20251210141230.jpg" # 您的拉伸机照片
        ],
        
        # --- Slide 6: Fazit (总结) ---
        [
            "Zusammenfassung",
            "Das untersuchte Pflaster besteht aus drei funktionellen Schichten:\n\n"
            "1. Träger: Perforiertes Polyethylen (PE) für Stabilität und Flexibilität.\n"
            "2. Wundauflage: Saugfähiges Zellulose-Vlies.\n"
            "3. Klebstoff: Haftstarkes Acrylat-Copolymer.\n\n"
            "Mechanisch zeichnet sich das Produkt durch hohe Duktilität aus, was die Anpassungsfähigkeit an Körperbewegungen gewährleistet.",
            None # 无图
        ]
    ]

    # --- 循环生成 ---
    for title_text, content_text, img_name in slides_data:
        slide_layout = prs.slide_layouts[1] # Content with Caption layout
        slide = prs.slides.add_slide(slide_layout)
        
        # 标题
        slide.shapes.title.text = title_text
        
        # 正文
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content_text
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18) # 调整字体大小
            paragraph.space_after = Pt(10)

        # 图片处理
        if img_name:
            img_path = os.path.join(image_folder, img_name)
            if os.path.exists(img_path):
                # 图片位置参数 (左, 上, 高度)
                left = Inches(5.5) 
                top = Inches(1.5)
                height = Inches(4.5)
                try:
                    slide.shapes.add_picture(img_path, left, top, height=height)
                except Exception as e:
                    print(f"Error loading {img_name}: {e}")

    prs.save(output_filename)
    print(f"PPT gespeichert: {output_filename}")

# --- 运行部分 ---
# 简单运行：图片与脚本在同一目录时，直接使用脚本所在目录。
if __name__ == "__main__":
    # 脚本所在目录（绝对路径）
    folder_path = os.path.abspath(os.path.dirname(__file__))

    print(f"使用图片文件夹: {folder_path}")
    create_german_lab_ppt_no_spectra(folder_path)